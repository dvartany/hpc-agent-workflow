#!/usr/bin/env python3
"""ASTRA slide: screenshot left, bulleted discussion right. Tight layout."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SCREENSHOT = "/Users/davidvartanyan/Desktop/Screenshot 2026-06-09 at 1.55.11 PM.png"
OUTPUT = "HPC_Agent_Slide_Deck.pptx"

IMG_W, IMG_H = Image.open(SCREENSHOT).size
ASPECT = IMG_W / IMG_H

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1B, 0x2A, 0x3D)
SLATE = RGBColor(0x3D, 0x52, 0x66)
BODY = RGBColor(0x4A, 0x5A, 0x6A)
RED = RGBColor(0xD9, 0x4A, 0x3C)
GREEN = RGBColor(0x1E, 0x9B, 0x56)
AMBER = RGBColor(0xE6, 0x8A, 0x0A)
BLUE = RGBColor(0x00, 0x7A, 0xCC)
SOFT_GRAY = RGBColor(0xE8, 0xEC, 0xF0)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

# ─── HEADER ───
hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.6))
hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK; hdr.line.fill.background()
tf = hdr.text_frame; tf.margin_left = Inches(0.4); tf.word_wrap = False
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "ASTRA  "; r1.font.size = Pt(22); r1.font.color.rgb = BLUE; r1.font.bold = True
r2 = p.add_run(); r2.text = "Automated Scheduling & Tracking for Research Applications"
r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)

# ─── SCREENSHOT — vertically centered on left ───
avail_h = 7.1 - 0.75  # 6.35 from below header to bottom bar
pic_w = Inches(7.8)
pic_h = Inches(7.8 / ASPECT)
pic_left = Inches(0.3)
pic_top = Inches(0.75 + (avail_h - pic_h.inches) / 2)
slide.shapes.add_picture(SCREENSHOT, pic_left, pic_top, pic_w, pic_h)

# Caption below screenshot
cap_y = pic_top.inches + pic_h.inches + 0.06
cap_h = Inches(0.65)
tx = slide.shapes.add_textbox(pic_left, Inches(cap_y), pic_w, cap_h)
tx.text_frame.word_wrap = True
p = tx.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = "Dashboard"; r1.font.size = Pt(10); r1.font.color.rgb = DARK; r1.font.bold = True
r2 = p.add_run(); r2.text = ": controls, config, logs"; r2.font.size = Pt(10); r2.font.color.rgb = SLATE
p.alignment = PP_ALIGN.LEFT

body_lines = [
    "Left panel: one-click pipeline buttons (Start Monitor → Pre-process → Submit → Analyze → Sync), phase indicators,",
    "config grid for cluster, jobs, analysis modes. Right panel: HPC connection settings, dynamic job scripts,",
    "compact status + run logs. All changes saved to config.toml automatically.",
]
tx2 = slide.shapes.add_textbox(Inches(0.35), Inches(cap_y + 0.18), pic_w, Inches(0.55))
tx2.text_frame.word_wrap = True
for i, line in enumerate(body_lines):
    if i == 0:
        p = tx2.text_frame.paragraphs[0]
    else:
        p = tx2.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(8); p.font.color.rgb = SLATE
    p.space_after = Pt(1)

# ─── RIGHT COLUMN ───
RX = 8.35
RW = 4.7

def add_section(y_top, color, number, title, summary, details):
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(RX), Inches(y_top), Inches(0.5), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()

    # Numbered title
    tx = slide.shapes.add_textbox(Inches(RX), Inches(y_top + 0.08), Inches(RW), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = f"{number}. {title}"
    p.font.size = Pt(15); p.font.color.rgb = DARK; p.font.bold = True

    # Summary sentence
    tx1 = slide.shapes.add_textbox(Inches(RX), Inches(y_top + 0.36), Inches(RW), Inches(0.35))
    tx1.text_frame.word_wrap = True
    p = tx1.text_frame.paragraphs[0]
    p.text = summary
    p.font.size = Pt(12); p.font.color.rgb = DARK; p.font.bold = True
    p.space_after = Pt(2)

    # Detail bullets
    tx2 = slide.shapes.add_textbox(Inches(RX), Inches(y_top + 0.65), Inches(RW), Inches(1.4))
    tx2.text_frame.word_wrap = True
    for i, d in enumerate(details):
        if i == 0:
            p = tx2.text_frame.paragraphs[0]
        else:
            p = tx2.text_frame.add_paragraph()
        p.text = f"  •  {d}"
        p.font.size = Pt(11.5); p.font.color.rgb = BODY
        p.space_after = Pt(1)


add_section(0.80, RED, 1, "Motivation",
    "HPC research today means juggling SSH terminals, sbatch scripts, and manual file transfers — fragile and error-prone.",
    [
        "Disconnected pre-process, submit, analysis, and sync steps",
        "No unified status view → wasted compute hours",
        "Goal: one repeatable workflow any lab member can run",
    ])

add_section(2.60, GREEN, 2, "Process",
    "One dashboard runs the full pipeline: pre-process → submit → poll → analyze → sync — across every job script.",
    [
        "Start Monitor iterates all job scripts in their own subdirectory",
        "Phase indicators light up for each pipeline step in real time",
        "Stop buttons scancel the Slurm job, then kill the local process",
        "Settings persist in config.toml — reproducible every time",
    ])

add_section(4.50, AMBER, 3, "Key Challenges",
    "Multi-job HPC pipelines break in predictable ways without the right safeguards built in.",
    [
        "Reproducibility — each script runs in its own remote workspace",
        "Scalability — add job scripts without touching the workflow",
        "Transparency — live logs so you never SSH blind again",
        "Safety — cancel at the scheduler level, not just locally",
    ])

# Closing line
tx = slide.shapes.add_textbox(Inches(RX), Inches(6.15), Inches(RW), Inches(0.45))
tx.text_frame.word_wrap = True
p = tx.text_frame.paragraphs[0]
p.text = "Stop babysitting individual batch jobs — automate the pipeline and focus on interpreting results."
p.font.size = Pt(12); p.font.color.rgb = DARK; p.font.bold = True

# ─── BOTTOM ───
bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
bot.fill.solid(); bot.fill.fore_color.rgb = DARK; bot.line.fill.background()
tf = bot.text_frame; tf.margin_left = Inches(0.4)
p = tf.paragraphs[0]
p.text = "ASTRA — github.com/dvartany/hpc-agent-workflow     ·     Built with opencode agents"
p.font.size = Pt(10); p.font.color.rgb = RGBColor(0x99, 0xAA, 0xBB)

prs.save(OUTPUT)
print(f"Saved {OUTPUT} ({len(prs.slides)} slide)")
print(f"Screenshot: {pic_w.inches:.1f}×{pic_h.inches:.1f} at y={pic_top.inches:.2f}")
